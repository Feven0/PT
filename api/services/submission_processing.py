import abc
import os, sys
import io
import re
import json
import base64
import zipfile
import pandas as pd
import tempfile
import shutil
import boto3
import requests
from bs4 import BeautifulSoup
import pypdf
import docx2txt
import docx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from striprtf.striprtf import rtf_to_text
import fitz
import openpyxl
import trafilatura
from PIL import Image
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload
import pandas as pd
from tabulate import tabulate
from kedro.io import AbstractDataset

from api.core import config
import src.llpacker.utils as sutils
import src.llpacker.services.secret as secret
from src.llpacker.services.llpacker_logger import LLPackerLogger
from api.services.strapi_graphql import StrapiGraphql
from src.tenx_auto_grade.data_sets.weaviate_dataset import WeaviateDataset
from src.llpacker.services.code_extractor import GitCodeAnalysis
from google.auth.exceptions import RefreshError
from src.llpacker.services.evaluation_error_ import *

logger = LLPackerLogger(os.path.basename(__file__))

def fetch_url(url, headers=None, redirect=True, content_type='text/html'):
    if headers is None:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) \
                            AppleWebKit/537.36 (KHTML, like Gecko) \
                            Chrome/90.0.4430.212 Safari/537.36',
            'Content-Type': content_type
        }
    return requests.get(url, headers=headers, allow_redirects=redirect)

def get_encoded_image_from_url(url, encode=False):
    if not encode:
        if url.startswith('http'):
            return url
        else:
            return None
    
    try:
        response = fetch_url(url)
        if response.status_code == 200:
             #if image size is too small (<10kb) or too large (10mb), return None
            if len(content := response.content) < 10000 or len(content) > 10000000: 
                img = None
            else:            
                img = base64.b64encode(response.content).decode('utf-8')

            return img
        else:
            logger.error(f"Unable to fetch image from url={url} \n status={response.status_code}")
            return None
    except Exception as e:
        logger.error(f"Unable to fetch image from url={url} \n {e}")
        return None

def get_encoded_image_from_url_list(url_list):
    try:
        encoded_images = []
        for url in url_list:
            img = get_encoded_image_from_url(url)
            if img is not None:
                encoded_images.append(img)
        return encoded_images
    except Exception as e:
        logger.error(f"Unable to encode image from url_list={url_list} \n {e}")
        return []
def extract_text_and_base64_images(file_io):
        with fitz.open(stream=file_io, filetype="pdf") as doc:
            text = ""
            encoded_images = []

            for page_num in range(len(doc)):
                page = doc[page_num]
                text += page.get_text()

                image_list = page.get_images(full=True)
                for img in image_list:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    encoded_image = base64.b64encode(image_bytes).decode('utf-8')
                    encoded_images.append(encoded_image)

        return text, encoded_images
def extract_text_and_images(file_bytes):
        with io.BytesIO(file_bytes) as f:
            pdf_reader = pypdf.PdfReader(f)
            text = ''
            images = []
            for page in pdf_reader.pages:
                text += page.extract_text()
                images.extend(page.images)
            return text, images
class DocumentProcessing():
    def __init__(self, **kwargs) -> None:        
        self.submission_id_col = kwargs.get("submission_id_col", "submission_id")
        self.url_col = kwargs.get("url_col", "url")
        self.url_type_col = kwargs.get("url_type_col", "url_type")
        self.doc_ext_col = kwargs.get("doc_ext_col", "document_format")
        self.doc_type_col = kwargs.get("doc_type_col", "document_type")
        self.doc_title_col = kwargs.get("doc_title_col", "document_title") 
        self.change_status = kwargs.get("change_status", ['Modified'])
        self.is_openai_file_id = kwargs.get("is_openai_file_id", False)
        self.run_stage = kwargs.get("run_stage", "dev")
        for _ in range(5):
            self.service_account_file = secret.get_google_service_account()        
            if self.service_account_file is None:
                logger.error("Google Service Account file not provided - it is None")  
                        
            try:
                credentials = service_account.Credentials.from_service_account_info(self.service_account_file)
                break
            except Exception as e:
                logger.error(f"Error in reading google service account file: {e}")
                continue

        self.drive_service = build('drive', 'v3', credentials=credentials, cache_discovery=False)   
        self.slides_service = build('slides', 'v1', credentials=credentials)
        self.docs_service = build('docs', 'v1', credentials=credentials)
        self.sheet_service = build('sheets', 'v4', credentials=credentials)

        self.since = kwargs.get("since", None)
        self.until = kwargs.get("until", None)
        self.skip_others_contribution = kwargs.get("skip_others_contribution", True)

        self.temp_dir = '/tmp'
    def __enter__(self):
        self.temp_dir = tempfile.mkdtemp()
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        if self.temp_dir:
            shutil.rmtree(self.temp_dir)
    def get_logfile_name(self, d):
        if isinstance(d, dict):
            submission_id = d[self.submission_id_col]
        elif isinstance(d, str):
            submission_id = d
        else:
            return f"submission_document_loader.log"
        return f"submission_document_loader_{submission_id}.log"
    
    def google_file_id_from_url(self, url):        
        if 'google.com' in url and '/d/' in url:            
            t = re.search(r'[-\w]{25,}(?!.*[-\w]{25,})', url)
            if t:
                t = t.group(0)
            else:
                t = ""
            
            if len(t) <= 25:
                t = url.split('/d/')[1].split('/')[0]

            if len(t) >= 25:
                return t
            else:
                logger.error(f'--Unable to extract google file id from passed url={url}')
                return url
        else:
            return url
        
    def is_valid_google_drive_id(self, text):
        drive_id_pattern = re.compile(r'^[a-zA-Z0-9_-]{25,}$')
        return bool(drive_id_pattern.match(text))
    def download_file_from_drive(self, file_id, bcontent=False):
        try:
            request = self.drive_service.files().get_media(fileId=file_id)
            file = io.BytesIO()
            downloader = MediaIoBaseDownload(file, request)
            done = False
            while done is False:
                status, done = downloader.next_chunk()
                
            if bcontent:
                return file.getvalue()
            else:
                return file
        except Exception as e:
            logger.error(f"Error in download_file_from_drive(file_id={file_id})")
            raise
    def dataframe_to_text_table(self, df):
            return tabulate(df,tablefmt='grid', showindex=False)

    # ******** Extract text and images from PDF ********
    def extract_pdf_content(self, file_io):
        if not self.temp_dir:
            self.temp_dir = tempfile.mkdtemp()

        with fitz.open(stream=file_io, filetype="pdf") as doc:
            full_text = ""
            images = {}
            image_index = 0

            for page_num, page in enumerate(doc):
                # Extract tables first
                tables = page.find_tables()
                table_areas = [table.bbox for table in tables]
                
                # Extract text
                text = page.get_text()
                
                # Process tables and insert them into the text
                for i, table_rect in enumerate(table_areas):
                    table_text = page.get_text("text", clip=table_rect)
                    df = pd.DataFrame(tables[i].extract())
                    formatted_table = self.dataframe_to_text_table(df)
                    text = text.replace(table_text, f"\n{formatted_table}\n")

                full_text += text + "\n"

                # Process images
                for img_index, img in enumerate(page.get_images(full=True)):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]

                    # Generate a unique image ID
                    image_id = f"IMAGE_ATTACHMENT_{page_num}_{img_index}"
                    image_filename = f"{image_id}.png"
                    image_path = os.path.join(self.temp_dir, image_filename)

                    # Save the image
                    pil_img = Image.open(io.BytesIO(image_bytes))
                    pil_img.save(image_path)

                    images[image_id] = {
                        "path": image_path
                    }

                    # Insert image placeholder in the text
                    full_text += f"[{image_id}]\n"
                    image_index += 1

        return full_text, images
    # ******** Method to handle native Google Doc ********
    def get_text_and_images_from_google_doc(self, file_id):
        try:
            doc = self.docs_service.documents().get(documentId=file_id).execute()
            doc_content = doc.get('body', {}).get('content', [])
            if not doc_content:
                logger.warning(f"No content found in Google Doc with ID: {file_id}")
                return "", {}

            def read_gdoc_paragraph_element(element):
                text_run = element.get('textRun')
                if not text_run:
                    return ''
                return text_run.get('content')

            def read_gdoc_structural_elements(elements):
                text = ''
                images = {}
                image_index = 0
                for value in elements:
                    if 'paragraph' in value:
                        elements = value.get('paragraph').get('elements')
                        for elem in elements:
                            text += read_gdoc_paragraph_element(elem)
                    elif 'table' in value:
                        table = value.get('table')
                        for row in table.get('tableRows'):
                            cells = row.get('tableCells')
                            for cell in cells:
                                cell_text, cell_images = read_gdoc_structural_elements(cell.get('content'))
                                text += cell_text
                                images.update(cell_images)
                    elif 'tableOfContents' in value:
                        toc = value.get('tableOfContents')
                        toc_text, toc_images = read_gdoc_structural_elements(toc.get('content'))
                        text += toc_text
                        images.update(toc_images)
                    elif 'inlineObjects' in value:
                        inline_objects = value.get('inlineObjects')
                        for obj_id, obj in inline_objects.items():
                            image_properties = obj.get('inlineObjectProperties', {}).get('embeddedObject', {}).get('imageProperties', {})
                            if 'contentUri' in image_properties:
                                image_url = image_properties['contentUri']
                            elif 'sourceUri' in image_properties:
                                image_url = image_properties['sourceUri']
                            else:
                                continue

                            try:
                                response = requests.get(image_url)
                                if response.status_code == 200:
                                    image_id = f"IMAGE_ATTACHMENT_{image_index}"
                                    image_filename = f"{image_id}.png"
                                    image_path = os.path.join(self.temp_dir, image_filename)

                                    with open(image_path, 'wb') as img_file:
                                        img_file.write(response.content)

                                    images[image_id] = {"path": image_path}
                                    text += f"[{image_id}]\n"
                                    image_index += 1
                            except Exception as e:
                                logger.error(f"Error downloading image from Google Doc: {str(e)}")

                return text, images

            text, images = read_gdoc_structural_elements(doc_content)
            if not text.strip():
                logger.warning("No text extracted from Google Doc. Falling back to PDF export.")
                return self.export_google_document(file_id)
            return text, images
        except Exception as e:
            logger.error(f"Error extracting content from Google Doc (ID: {file_id}): {str(e)}")
            logger.info("Attempting to extract content from PDF export.")
            return self.export_google_document(file_id)   
    # ******** Method to handle native Google Sheets ********
    def read_google_sheets(self, file_id, max_rows=100, max_cols=50):
        logger.info("read_google_sheets")
        try:
            sheet_metadata = self.sheet_service.spreadsheets().get(spreadsheetId=file_id).execute()
            sheets = sheet_metadata.get('sheets', [])

            full_text = ""

            for sheet in sheets:
                sheet_name = sheet['properties']['title']
                full_text += f"\n\n--- Sheet: {sheet_name} ---\n\n"

                # Get sheet dimensions
                grid_properties = sheet['properties']['gridProperties']
                sheet_rows = grid_properties.get('rowCount', 0)
                sheet_cols = grid_properties.get('columnCount', 0)

                # Define range to fetch
                range_name = f"'{sheet_name}'!A1:{self.column_letter(max_cols)}{max_rows}"
                result = self.sheet_service.spreadsheets().values().get(
                    spreadsheetId=file_id, range=range_name).execute()
                values = result.get('values', [])

                if not values:
                    full_text += "No data found in this sheet.\n"
                    continue

                df = pd.DataFrame(values[1:], columns=values[0])
                full_text += self.dataframe_to_text_table(df) + "\n"

                if sheet_rows > max_rows or sheet_cols > max_cols:
                    full_text += f"\nNote: This sheet has been truncated. Original size: {sheet_rows} rows, {sheet_cols} columns.\n"

            return full_text, {}  # Return an empty dict for images
        except Exception as e:
            logger.error(f"Unable to get text from Google Sheets file id={file_id}: {e}")
            return "", {}

    def column_letter(self, column_number):
        """Convert a column number to a column letter (A, B, C, ..., Z, AA, AB, ...)"""
        column_letter = ''
        while column_number > 0:
            column_number, remainder = divmod(column_number - 1, 26)
            column_letter = chr(65 + remainder) + column_letter
        return column_letter

    # ******** Method to handle native Google Slides ********
    def get_text_and_images_from_google_slide(self, file_id):
        try:
            text, images = self.read_google_slides(file_id)
            if not text.strip():
                logger.error("No text extracted from Google Slides. Falling back to PDF export.")
                return self.export_google_document(file_id)
            return text, images
        except Exception as e:
            logger.error(f"Error extracting content from Google Slides (ID: {file_id}): {str(e)}")
            logger.info("Attempting to extract content from PDF export.")
            return self.export_google_document(file_id)
    def read_google_slides(self, presentation_id):
        def get_slide_text_element(element):
            text_list = []
            if 'shape' in element and 'text' in element['shape']:
                text_elements = element['shape']['text']['textElements']
                for text_element in text_elements:
                    if 'textRun' in text_element and 'content' in text_element['textRun']:
                        text = text_element['textRun']['content']
                        text_list.append(text)
            return '\n'.join(text_list)

        def get_slide_image_element(element):
            image_list = []
            if 'image' in element:
                if isinstance(element['image'], dict):
                    images = [element['image']]
                elif isinstance(element['image'], list):
                    images = element['image']
                else:
                    images = []

                for image_item in images:
                    img_url = ""
                    if 'contentUrl' in image_item:
                        img_url = image_item['contentUrl']
                    elif 'sourceUrl' in image_item:
                        img_url = image_item['sourceUrl']

                    if img_url:
                        image_content = self.fetch_url(img_url).content
                        image_id = f"IMAGE_ATTACHMENT_{len(image_list)}"
                        image_filename = f"{image_id}.png"
                        image_path = os.path.join(self.temp_dir, image_filename)

                        with open(image_path, 'wb') as img_file:
                            img_file.write(image_content)

                        image_list.append({image_id: {"path": image_path}})
            return image_list

        def get_slide_table_element(element):
            if 'table' in element:
                table_data = []
                for row in element['table']["tableRows"]:
                    row_data = []
                    for cell in row["tableCells"]:
                        cell_text = ""
                        for text_element in cell["text"]["textElements"]:
                            if 'textRun' in text_element and 'content' in text_element['textRun']:
                                cell_text += text_element['textRun']['content']
                        row_data.append(cell_text)
                    table_data.append(row_data)
                
                df = pd.DataFrame(table_data)
                return self.dataframe_to_text_table(df)
            return ""

        credentials = service_account.Credentials.from_service_account_info(self.service_account_file)
        slides_service = build('slides', 'v1', credentials=credentials)

        try:
            presentation = slides_service.presentations().get(presentationId=presentation_id).execute()
        except Exception as e:
            logger.error('Error in slides_service.presentations().get')
            raise

        full_text = f"Slide title: {presentation.get('title', '')}\n"
        full_text += f"Slide description: {presentation.get('description', '')}\n\n"

        images = {}
        image_index = 0

        for slide_index, slide in enumerate(presentation['slides']):
            full_text += f"Slide {slide_index + 1}:\n"
            full_text += f"Slide title: {slide.get('title', '')}\n"

            for element in slide['pageElements']:
                text_element = get_slide_text_element(element)
                table_element = get_slide_table_element(element)
                image_elements = get_slide_image_element(element)

                if text_element:
                    full_text += text_element + "\n"

                if table_element:
                    full_text += f"\n{table_element}\n"

                for image_element in image_elements:
                    for image_id, image_info in image_element.items():
                        images[image_id] = image_info
                        full_text += f"[{image_id}]\n"
                        image_index += 1

            full_text += "\n"  # Add a blank line between slides

        return full_text, images
    # *********** Method to handle native Google PDF ***********    
    def read_pdf_from_gdrive(self, file_id):
        logger.info(f"Calling read_pdf_from_gdrive(file_id={file_id}) ...")

        if not self.is_valid_google_drive_id(file_id):
            raise ValueError("Invalid Google Drive file ID")

        file_io = self.download_file_from_drive(file_id, bcontent=False)

        try:            
            text, images = self.extract_pdf_content(file_io)                 
            return text, images
        except Exception as e:
            logger.error(f"Error in calling extract_text_and_images: {e}")
            raise
    # ******** Export as pdf from google slides and google doc ********
    def export_google_document(self, file_id, export_format='application/pdf'):
        request = self.drive_service.files().export(fileId=file_id, 
                                                    mimeType=export_format)
        file_io = io.BytesIO()
        downloader = MediaIoBaseDownload(file_io, request)
        done = False
        while done is False:
            status, done = downloader.next_chunk()
        try:
            if export_format == 'application/pdf':
                text, images = self.extract_pdf_content(file_io)         
                text = text.strip()
            elif export_format == 'text/plain':
                file_bytes = file_io.getvalue()
                text = file_bytes.decode('utf-8')
                images = []
            else:
                file_bytes = file_io.getvalue()
                text = file_bytes.decode('utf-8')
                images = []
            return text, images
        except Exception as e:           
            logger.error(f"Error in exporting google drive file: file_id={file_id}, export_format={export_format})")
            raise           

    # *********** Method to handle Excel files stored in Google Drive. ***********    
    def read_google_drive_excel_file(self, file_id, max_rows=1000, max_cols=50):
        try:
            request = self.drive_service.files().get_media(fileId=file_id)
            excel_content = request.execute()
            excel_file = io.BytesIO(excel_content)
            
            workbook = openpyxl.load_workbook(excel_file, data_only=True, read_only=True)
            full_text = ""

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                full_text += f"\n\n--- Sheet: {sheet_name} ---\n\n"

                # Get dimensions of the sheet
                max_row = min(sheet.max_row, max_rows)
                max_col = min(sheet.max_column, max_cols)

                # Extract limited data
                data = []
                for row in sheet.iter_rows(max_row=max_row, max_col=max_col, values_only=True):
                    data.append(row)

                if not data:
                    full_text += "No data found in this sheet.\n"
                    continue

                df = pd.DataFrame(data[1:], columns=data[0])
                full_text += self.dataframe_to_text_table(df) + "\n"

                if sheet.max_row > max_rows or sheet.max_column > max_cols:
                    full_text += f"\nNote: This sheet has been truncated. Original size: {sheet.max_row} rows, {sheet.max_column} columns.\n"

            return full_text, {}  # Return an empty dict for images
        except Exception as e:
            logger.error(f"Unable to get text from Google Drive Excel file id={file_id}: {e}")
            return "", {}


    # *********** Method to handle Word files stored in Google Drive.***********
    def read_google_drive_word_file(self, file_id):
        try:
            if not self.is_valid_google_drive_id(file_id):
                raise InvalidInputError("Invalid Google Drive file ID")
            
            request = self.drive_service.files().get_media(fileId=file_id)
            file_content = io.BytesIO(request.execute())

            document = docx.Document(file_content)
            full_text = ""
            images = {}
            image_index = 0

            for paragraph in document.paragraphs:
                full_text += paragraph.text + "\n"

            for rel in document.part.rels.values():
                if "image" in rel.target_ref:
                    image_bytes = rel.target_part.blob
                    image_id = f"IMAGE_ATTACHMENT_{image_index}"
                    image_filename = f"{image_id}.png"
                    image_path = os.path.join(self.temp_dir, image_filename)

                    with open(image_path, 'wb') as img_file:
                        img_file.write(image_bytes)

                    images[image_id] = {"path": image_path}
                    full_text += f"[{image_id}]\n"
                    image_index += 1

            return full_text, images
        except Exception as e:
            logger.error(f"Unable to get text from Google Drive Word file id={file_id}: {e}")
            return "", {}
        
    # *********** Method to handle PPT files stored in Google Drive. ***********
    def read_google_drive_pptx(self, file_id):
        try:
            request = self.drive_service.files().get_media(fileId=file_id)
            pptx_contents = request.execute()
            pptx_stream = io.BytesIO(pptx_contents)
            slide_text, images = self.read_pptx_from_bytesio(pptx_stream)
            return slide_text, images 
        except Exception as e:
            logger.error(f"Unable to get text from google drive pptx file id={file_id}: {e}")
            return "", []
    def read_pptx_from_bytesio(self, bytes_io):
        pptx = Presentation(bytes_io)
        full_text = ""
        images = {}
        image_index = 0

        for slide_index, slide in enumerate(pptx.slides):
            full_text += f"Slide {slide_index + 1}:\n"

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        full_text += paragraph.text + "\n"
                
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    df = pd.DataFrame(table_data)
                    full_text += f"\n{self.dataframe_to_text_table(df)}\n"

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob
                    image_id = f"IMAGE_ATTACHMENT_{image_index}"
                    image_filename = f"{image_id}.png"
                    image_path = os.path.join(self.temp_dir, image_filename)

                    with open(image_path, 'wb') as img_file:
                        img_file.write(image_bytes)

                    images[image_id] = {"path": image_path}
                    full_text += f"[{image_id}]\n"
                    image_index += 1

            full_text += "\n"  # Add a blank line between slides

        return full_text, images
    def select_gdrive_reader(self, ext, exact=False):
        reader_func = {
            'pdf': {"name": "read_pdf_from_gdrive", "func": self.read_pdf_from_gdrive},
            'gdoc': {"name": "read_google_doc_as_word_file", "func": self.get_text_and_images_from_google_doc},
            'docx': {"name": "read_google_drive_word_file", "func": self.read_google_drive_word_file},
            "pptx": {"name": "read_google_drive_pptx", "func": self.read_google_drive_pptx},
            "gslide": {"name": "read_google_slides", "func": self.read_google_slides},
            "excel": {"name": "read_google_drive_excel", "func": self.read_google_drive_excel_file},
            "gsheet": {"name": "read_google_sheets", "func": self.read_google_sheets},
            'default': {"name": "read_pdf_from_gdrive", "func": self.read_pdf_from_gdrive}
        }
        reader_func['application/vnd.google-apps.presentation'] = reader_func['gslide']
        reader_func['application/vnd.google-apps.document'] = reader_func['gdoc']
        reader_func['application/vnd.google-apps.spreadsheet'] = reader_func['gsheet']
        reader_func['application/pdf'] = reader_func['pdf']
        reader_func['application/vnd.openxmlformats-officedocument.wordprocessingml.document'] = reader_func['docx']
        reader_func['application/vnd.openxmlformats-officedocument.presentationml.presentation'] = reader_func['pptx']
        reader_func['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'] = reader_func['excel']

        if exact:
            if ext in reader_func.keys():
                return reader_func[ext]
            else:
                return None
        else:
            return reader_func.get(ext, reader_func['default'])
    # *********** Associate google drive file with reader ***********
    def read_gdrive_file(self, file_id, file_name=None):
        try:
            request = (self.drive_service.files()
                       .get(fileId=file_id, fields='*')
                       .execute()
            )

            mimetype = request['mimeType']
            print(f"mimetype={mimetype}")
            file_name = request['name']
            file_id = request['id']

            logger.info(f"Calling select_gdrive_reader with {mimetype}")
            reader = self.select_gdrive_reader(mimetype, exact=True)

            skip = 'vnd.google-apps' in mimetype
            if reader is None or skip:
                if mimetype == 'application/vnd.google-apps.presentation':
                    logger.info(f"Extracting text and images from google slides ...")
                    try:
                        text, images = self.get_text_and_images_from_google_slide(file_id)
                    except Exception as e:
                        logger.error(f"Error in calling get_text_and_images_from_google_slide: {e}")
                        raise

                elif mimetype == 'application/vnd.google-apps.document':
                    logger.info(f"Extracting text and images from google document ...")
                    try:
                        text, images = self.get_text_and_images_from_google_doc(file_id)
                    except Exception as e:
                        logger.error(f"Error in calling get_text_and_images_from_google_doc: {e}")
                        raise

                elif mimetype == 'application/vnd.google-apps.spreadsheet':
                    text, images = self.read_google_sheets(file_id)
                else:
                    text, images = "", []              
            else:
                logger.info(f"Calling {reader['name']} with file_id={file_id} ...")
                text, images = reader['func'](file_id)

            return text, images
        except InvalidInputError as e:
            logger.error(f"Unable to get text from google doc file id={file_id}: {e}")
            raise 
    def check_for_url_validation(self, url):
        url_pattern = re.compile(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', re.IGNORECASE)
        return re.match(url_pattern, url) is not None

    def fetch_images_from_html_document(self, html_document, encoded=False):
        soup = BeautifulSoup(html_document, 'html.parser')
        imgtag = "img"
        for tag in ['figure','picture','img']:
            figure_tags = soup.find_all(tag)
            if len(figure_tags)>0:
                imgtag = tag
                break

        images = []
        for item in figure_tags:            
            try:
                if imgtag in ['figure', 'picture']:
                    image_url = item.find('source')['srcset'].split(',')[-1].strip().split(' ')[0]
                else:
                    image_url = item['src'].strip()

                if len(image_url)==0:
                    continue
            except:
                continue

            if encoded:
                img = get_encoded_image_from_url(image_url)
            else:
                img = image_url

            if img is not None:
                try:
                    caption = item.find('figcaption').text
                except:
                    caption = ""

                images.append((img, caption))

        return images

    def read_text_from_link(self, url):
        if url == 'https://medium.com/p/ceaf6db24b7f/edit':
            url = "https://medium.com/@eyayab21/ceaf6db24b7f"

        try:
            if self.check_for_url_validation(url):
                headers = {
                        'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/90.0.4430.212 Safari/537.36'
                    }
                    
                html_document = fetch_url(url).content                
                article_trafilatura = trafilatura.bare_extraction(html_document)
                if article_trafilatura is None:
                    logger.warn(f"Unable to extract text from url={url}") 
                    raise

                text = ""
                for k in ['title', 'url', 'date', 'text', 'description']:
                    if len(article_trafilatura.get(k, "")) > 0:
                        if k == 'text':
                            h = 'Body'
                        else:
                            h = k.title()

                        v = article_trafilatura[k]
                        if not v.replace("...","") in text:
                            text += f"#{h} - {v} \n"

                try:
                    images = self.fetch_images_from_html_document(html_document)
                except Exception as e:
                    logger.error(f"Unable to fetch images from html document {e}")
                    images = []

                return text, images
            else:
                raise InvalidInputError("Incorrect URL format")
        except Exception  as e:
            logger.error(f"Unable to read text from url={url}: {e}")
            raise

    def get_bucket_and_file_name(self, url):
        if url.startswith('s3://'):
            url = url.replace('s3://', '')
            bucket_name = url.split('/')[2]
            file_name = '/'.join(url.split('/')[3:])
            return bucket_name, file_name
        else:
            raise InvalidInputError("Incorrect S3 URL format")

    def read_from_s3(self, url):
        try:
            if 'amazonaws.com' in url and url.startswith('http'):
                return self.read_text_from_link(url)
            else:
                bucket_name, file_name = self.get_bucket_and_file_name(url)
                s3_client = boto3.client('s3')
                file_bytes = s3_client.get_object(Bucket=bucket_name, Key=file_name)['Body'].read()
                text = self.extract_text(file_bytes)
                text = text.strip()
                text = ''.join(text.split("\n"))
                return text, []
        except InvalidInputError as e:
            logger.error(f"Unable to read text from s3 file url={url}: {e}")
            return "", []
        
    def get_code_content(self, git_url, max_token=-1, **kwargs):
        logger.info(f"Loading codes from git_url={git_url} with Status = {self.change_status}")
        try:
            gca = GitCodeAnalysis(git_url, **kwargs)
            content = gca.parse_repo_content(**kwargs)
        except Exception as e:
            logger.error(f'Error in GitCodeAnalysis ... {e}')
            return ""
        
        if len(content) == 0:
            logger.warn(f'No content extracted git_url={git_url}')
            return ""
        else:
            logger.good(f"Successfully extracted content from git_url={git_url}")
            return content
                

    def process_single_document(self, rowX):
        try:
            if isinstance(rowX, dict):
                row = rowX
            elif isinstance(rowX, pd.Series):
                row = rowX.to_dict()
            else:
                raise InvalidTypeInputError("Invalid input for 'row'. Expected a pandas Series.")
        except Exception as e:
            logger.error('Error in process_single_document Input ... {e}')
            return "", []

        try:
            reader_map = {
                "link": {
                    'html': {"name": "read_text_from_link", "func": self.read_text_from_link},
                    'default': {"name": "read_text_from_link", "func": self.read_text_from_link}
                },
                "gdrive": {
                    'default': {"name": "read_gdrive_file", "func": self.read_gdrive_file}
                            },
                "s3": {
                    'pdf': {"name": "read_from_s3", "func": self.read_from_s3},
                    'default': {"name": "read_from_s3", "func": self.read_from_s3}
                },
                "git": {
                    'default': {"name": "get_code_content", "func": self.get_code_content}                            
                },
                "text": {
                    'txt': {"name": "parse_simple_text", "func": lambda x: (x, [])},
                    'text': {"name": "parse_simple_text", "func": lambda x: (x, [])},
                    'default': {"name": "parse_simple_text", "func": lambda x: (rtf_to_text(x, errors="ignore"), [])}
                }
            }

            for x in ['git', 'gdrive', 's3']:
                reader_map['link'].update(reader_map[x])

            link_alias = ['url', 'link']
            for x in link_alias:
                if x != 'link':
                    reader_map[x] = reader_map['link']

            gdrive_alias = ['gdrive', 'google_drive', 'google']
            for x in gdrive_alias:
                if x != 'gdrive':
                    reader_map[x] = reader_map['gdrive']

            s3_alias = ['s3', 'aws']
            for x in s3_alias:
                if x != 's3':
                    reader_map[x] = reader_map['s3']

            git_alias = ['git', 'github', 'github-link', 'git-link', 'git_link', 'github_link']
            for x in git_alias:
                if x != 'git':
                    reader_map[x] = reader_map['git']

            url = row[self.url_col]
            url_type = row[self.url_type_col]
            ut, dt, df = sutils.get_url_type(url)
            print("url type=", ut, "documnet type=", dt,"documnet format=", df)
            if ut != url_type:
                logger.warn(f'found: url_type={url_type} but expected {ut}')
                url_type = ut

            ext = row[self.doc_ext_col]          
            urlToUse = url

            passed = f"url_type={url_type}, ext={ext}, url={url}"
            if len(url) == 0:
                logger.info(f'found: {passed} but expected len(url)>0 url_type in {reader_map.keys()}')
                raise InvalidInputError("Invalid url")    

            if url_type not in reader_map.keys():
                logger.warn(f'found: {passed} but expected url_type in {reader_map.keys()}')                
                url_type = ut
                ext = df

            if ext not in reader_map.get(url_type, {}).keys():
                ext = 'default'

            if url_type in gdrive_alias:
                urlToUse = self.google_file_id_from_url(url)
                print("urlToUse=", urlToUse)   
                print("url_type=", url_type)
                print("ext=", ext)             
                func = reader_map.get(url_type).get(ext, self.read_gdrive_file)
            elif url_type in s3_alias:
                urlToUse = url
                func = reader_map.get(url_type).get(ext, self.read_from_s3)
            elif url_type in git_alias:
                urlToUse = url
                func = reader_map.get(url_type).get(ext, self.get_code_content)
            else:
                urlToUse = url
                func = reader_map.get(url_type).get(ext, self.read_text_from_link)
        except Exception as e:
            logger.error('Error in selecting submission content reader function... {e}')
            return "", []
        try:    
            result = func["func"](urlToUse)
            if isinstance(result, tuple) and len(result) == 2:
                text, images = result
            elif isinstance(result, str):
                text, images = result, []
            else:
                logger.error(f'Unexpected return type from {func["name"]}')
                return "", []
        except Exception as e:  
            logger.error(f'Error in calling function={func["name"]}... {e}')    
            return "", []
            
        return text.strip(), images

    def extract_text_from_document(self, metadata):
        try:
            if isinstance(metadata, (dict, pd.Series)):
                text, images = self.process_single_document(metadata)
            elif isinstance(metadata, pd.DataFrame):
                text = []
                images = []
                for index, row in metadata.iterrows():
                    t, i = self.process_single_document(row)
                    text.append(t)
                    images.extend(i)
            elif isinstance(metadata, list):                
                text = []
                images = []
                for row in metadata:
                    t, i = self.process_single_document(row)
                    text.append(t)
                    images.extend(i)
            else:
                raise InvalidTypeInputError("Invalid input for 'metadata'. Expected a dict, pandas DataFrame, Series, or list.")
            return text, images
        except Exception as e:
            logger.error(f'Error in extract_text_from_document... {e}')
            return "", []

    def make_metadata(self, item):
        cols = [self.url_col,
                self.doc_title_col,
                self.doc_ext_col,
                self.doc_type_col]
        res = [f"{key}::{item.get(key)}" for key in item.keys() if key not in cols]
        return res
    
    def process_payloads(self, obj):        
        if 'payload' in obj.keys():
            payload = obj.get("payload")
        else:
            raise InvalidInputError("Invalid input. Expected a dictionary with payload or metadata key")

        params = obj.get("params", {})
        sid = obj.get("submission_id", "")        
        is_openai_file_id = params.get("is_openai_file_id", False)

        if isinstance(payload, str):
            payload = json.loads(payload)

        if not isinstance(payload, (list, tuple)):
            payload_list = [payload]
        else:
            payload_list = payload

        outputList = []
        #all_images = []

        for item in payload_list:
            if not isinstance(item, dict):     
                continue  

            if 'submission_id' not in item.keys():
                item['submission_id'] = sid

            url = item.get(self.url_col, "")

            text, images = self.extract_text_from_document(item)
            if text in ["", None] or len(text) < 10:
                logger.error(f"Unable to extract text from submission_id={sid} url={url}")
                continue
            else:
                logger.good(f"Successfully extracted text from submission_id={sid} url={url}")                
                outputList.append(text)
                #all_images.extend(images)
               
        return '\n'.join(outputList), images

    def extract_submission_content(self, payload, **kwargs) -> pd.DataFrame:
        logger.divider(f'Extracting submission content for sid={payload.get("submission_id")}')
        weaviate = kwargs.get("weaviate")

        weaviate_client_passed = isinstance(weaviate, AbstractDataset)
        
        if not weaviate_client_passed:
            logger.warn(f"Weaviate data catalog instance is not passed!")
            return {}        
        text, all_images = None, None    
        if weaviate_client_passed:
            text,payload = weaviate.get_submission_object(payload, run_stage=self.run_stage)

        if kwargs.get("force", False) or text in [None, "", {}, []]:
            if kwargs.get("force", False):
                logger.info(f"Extracting document as force=True ...")
            elif text in [None, "", {}, []]:
                logger.info(f"Extracting document as it is not in weaviate ...")

            text, all_images = self.process_payloads(payload)
            payload['images'] = all_images
            if text:           
                logger.good('preprocess_submission: text extracted from submission!')
                if weaviate_client_passed:
                    if all_images is not None:
                        payload['images'] = all_images
                        payload = weaviate.save_submission_object(text, payload, run_stage=self.run_stage)
            else:
                logger.error('preprocess_submission: text not extracted from submission!')
   
        return text,payload
